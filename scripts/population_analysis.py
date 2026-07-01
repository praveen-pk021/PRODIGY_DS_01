from pathlib import Path
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet

ROOT = Path(__file__).resolve().parents[1]
DATA_PATH = ROOT / "Dataset" / "world_population.csv"
IMAGES_DIR = ROOT / "Images"
REPORT_DIR = ROOT / "Report"
PRESENTATION_DIR = ROOT / "Presentation"
IMAGES_DIR.mkdir(exist_ok=True)
REPORT_DIR.mkdir(exist_ok=True)
PRESENTATION_DIR.mkdir(exist_ok=True)

plt.style.use("seaborn-v0_8-whitegrid")


def load_data(path: Path) -> pd.DataFrame:
    df = pd.read_csv(path)
    df.columns = [col.strip() for col in df.columns]

    for col in [
        "2022 Population",
        "2020 Population",
        "2015 Population",
        "2010 Population",
        "2000 Population",
        "1990 Population",
        "1980 Population",
        "1970 Population",
        "Area (km²)",
        "Density (per km²)",
        "Growth Rate",
    ]:
        df[col] = pd.to_numeric(df[col], errors="coerce")

    df["Continent"] = df["Continent"].fillna("Unknown").astype(str).str.strip()
    df = df.dropna(subset=["2022 Population"]).copy()
    return df


def compute_summary(df: pd.DataFrame) -> dict:
    pop = df["2022 Population"]
    return {
        "mean": round(float(pop.mean()), 2),
        "median": round(float(pop.median()), 2),
        "mode": round(float(pop.mode().iloc[0]), 2),
        "std": round(float(pop.std(ddof=0)), 2),
        "min": round(float(pop.min()), 2),
        "max": round(float(pop.max()), 2),
        "countries": int(df.shape[0]),
        "continents": int(df["Continent"].nunique()),
    }


def save_visualizations(df: pd.DataFrame) -> None:
    pop = df["2022 Population"]
    top_countries = df.nlargest(10, "2022 Population")
    continent_totals = df.groupby("Continent")["2022 Population"].sum().sort_values(ascending=False)

    fig, axes = plt.subplots(2, 2, figsize=(16, 12))
    fig.suptitle("Global Population Distribution Analysis", fontsize=18, fontweight="bold")

    sns.histplot(pop, bins=30, kde=True, color="#4C78A8", ax=axes[0, 0])
    axes[0, 0].set_title("Histogram with KDE for 2022 Population")
    axes[0, 0].set_xlabel("Population")
    axes[0, 0].set_ylabel("Number of Countries")

    sns.boxplot(x="Continent", y="2022 Population", data=df, ax=axes[0, 1], color="#4C78A8")
    axes[0, 1].set_title("Population Distribution by Continent")
    axes[0, 1].set_xlabel("Continent")
    axes[0, 1].set_ylabel("Population")
    axes[0, 1].tick_params(axis="x", rotation=45)

    sns.barplot(x="Country/Territory", y="2022 Population", data=top_countries, ax=axes[1, 0], color="#54A24B")
    axes[1, 0].set_title("Top 10 Most Populous Countries")
    axes[1, 0].set_xlabel("Country")
    axes[1, 0].set_ylabel("Population")
    axes[1, 0].tick_params(axis="x", rotation=45)

    sns.barplot(x=continent_totals.index, y=continent_totals.values, ax=axes[1, 1], color="#F58518")
    axes[1, 1].set_title("Total Population by Continent")
    axes[1, 1].set_xlabel("Continent")
    axes[1, 1].set_ylabel("Total Population")
    axes[1, 1].tick_params(axis="x", rotation=45)

    plt.tight_layout(rect=[0, 0, 1, 0.96])
    fig.savefig(IMAGES_DIR / "population_analysis_dashboard.png", dpi=300, bbox_inches="tight")
    plt.close(fig)

    # Save individual plots
    plt.figure(figsize=(10, 6))
    sns.histplot(pop, bins=30, kde=True, color="#4C78A8")
    plt.title("Population Distribution Histogram")
    plt.xlabel("Population")
    plt.ylabel("Number of Countries")
    plt.savefig(IMAGES_DIR / "population_histogram.png", dpi=300, bbox_inches="tight")
    plt.close()

    plt.figure(figsize=(10, 6))
    sns.kdeplot(pop, fill=True, color="#f58518")
    plt.title("Population Density Estimate")
    plt.xlabel("Population")
    plt.ylabel("Density")
    plt.savefig(IMAGES_DIR / "population_kde.png", dpi=300, bbox_inches="tight")
    plt.close()

    plt.figure(figsize=(10, 6))
    sns.barplot(x=top_countries["Country/Territory"], y=top_countries["2022 Population"], color="#54A24B")
    plt.title("Top 10 Countries by 2022 Population")
    plt.xlabel("Country")
    plt.ylabel("Population")
    plt.xticks(rotation=45, ha="right")
    plt.savefig(IMAGES_DIR / "top_countries_bar.png", dpi=300, bbox_inches="tight")
    plt.close()


def generate_report(df: pd.DataFrame, summary: dict) -> None:
    report_path = REPORT_DIR / "Task01_Report.pdf"
    title = "Population Distribution Analysis"
    story = []
    styles = getSampleStyleSheet()
    story.append(Paragraph(title, styles["Title"]))
    story.append(Spacer(1, 12))
    story.append(Paragraph("Prodigy InfoTech Data Science Internship - Task 01", styles["Heading2"]))
    story.append(Spacer(1, 12))

    summary_text = [
        f"This report analyzes the distribution of countries by population using the world_population.csv dataset.",
        f"The dataset contains {summary['countries']} countries and spans {summary['continents']} continents.",
        f"Key descriptive statistics for the 2022 population column are: mean = {summary['mean']}, median = {summary['median']}, mode = {summary['mode']}, standard deviation = {summary['std']}, minimum = {summary['min']}, and maximum = {summary['max']}.",
        "The histogram and KDE plot show that most countries have relatively small populations, while a small set of countries have very large populations.",
        "The bar chart of the top 10 countries highlights the dominance of populous nations such as India, China, and the United States.",
    ]
    for paragraph in summary_text:
        story.append(Paragraph(paragraph, styles["BodyText"]))
        story.append(Spacer(1, 8))

    top_countries = df.nlargest(5, "2022 Population")[["Country/Territory", "2022 Population"]]
    table_data = [["Country", "2022 Population"]]
    for _, row in top_countries.iterrows():
        table_data.append([row["Country/Territory"], f"{int(row['2022 Population']):,}"])

    table = Table(table_data, repeatRows=1)
    table.setStyle(
        TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4C78A8")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
            ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 10),
        ])
    )
    story.append(Spacer(1, 12))
    story.append(Paragraph("Top 5 Countries by Population", styles["Heading2"]))
    story.append(table)

    doc = SimpleDocTemplate(str(report_path), pagesize=letter)
    doc.build(story)


def generate_presentation(df: pd.DataFrame, summary: dict) -> None:
    presentation_path = PRESENTATION_DIR / "Task01_Presentation.pptx"
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Global Population Distribution Analysis"
    title_slide.placeholders[1].text = "Prodigy InfoTech | Task 01"

    slide_1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_1.shapes.title.text = "Objectives"
    body = slide_1.shapes.placeholders[1].text_frame
    body.text = "- Analyze the distribution of world population data\n- Visualize how population is spread across countries and continents\n- Highlight key descriptive statistics and notable countries"

    slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_2.shapes.title.text = "Key Findings"
    body = slide_2.shapes.placeholders[1].text_frame
    body.text = (
        f"- Countries analyzed: {summary['countries']}\n"
        f"- Continents represented: {summary['continents']}\n"
        f"- Mean population: {summary['mean']}\n"
        f"- Median population: {summary['median']}\n"
        f"- Highest populated country in the dataset: {df.nlargest(1, '2022 Population')['Country/Territory'].iloc[0]}"
    )

    slide_3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide_3.shapes.title.text = "Visual Highlights"
    text_box = slide_3.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(5.5), Inches(1.0))
    text_frame = text_box.text_frame
    text_frame.text = "The project includes a histogram, KDE plot, bar chart, and continent-level boxplot to communicate population spread clearly."
    slide_3.shapes.add_picture(str(IMAGES_DIR / "population_analysis_dashboard.png"), Inches(0.8), Inches(2.4), width=Inches(8.0))

    prs.save(presentation_path)


def main() -> None:
    df = load_data(DATA_PATH)
    summary = compute_summary(df)
    save_visualizations(df)
    generate_report(df, summary)
    generate_presentation(df, summary)

    print("Analysis completed successfully.")
    print(f"Summary: {summary}")


if __name__ == "__main__":
    main()
